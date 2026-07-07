import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import subprocess  # Handles the automatic execution of your gauge files

# -----------------------------
# 1. LOAD DATA FROM EXCEL
# -----------------------------
file_path = "clean_data.xlsx"

kpi_df = pd.read_excel(file_path, "KPIs")
sla_df = pd.read_excel(file_path, "SLA")
trend_df = pd.read_excel(file_path, "Ticket_Trend")
source_df = pd.read_excel(file_path, "Source_Availability")
issue_df = pd.read_excel(file_path, "Prod_Failure_Issue_Type")
batch_df = pd.read_excel(file_path, "Prod_Failure_Batch_Type")
dq_df = pd.read_excel(file_path, "DQ_Issues")

# Clean column names and text fields to remove unexpected padding
for df in [kpi_df, sla_df, trend_df, source_df, issue_df, batch_df, dq_df]:
    df.columns = df.columns.str.strip()
    for col in df.select_dtypes(include=['object']).columns:
        df[col] = df[col].astype(str).str.strip()

# -----------------------------
# 2. KPI + SLA MAPPING
# -----------------------------
kpi_map = dict(zip(kpi_df['Metric'], kpi_df['Value']))
sla_map = dict(zip(sla_df['Metric'], sla_df['Value']))

replacements = {
    "<<BATCH_CYCLES>>": str(kpi_map.get("BATCH_CYCLES", "")),
    "<<TOTAL_INCIDENTS>>": str(kpi_map.get("TOTAL_INCIDENTS", "")),
    "<<MTTR>>": str(kpi_map.get("MTTR", "")).strip(),
    "<<RECURRING_ISSUES>>": str(kpi_map.get("RECURRING_ISSUES", "")),
    "<<RECURRING_RESOLVED>>": str(kpi_map.get("RECURRING_RESOLVED", "")),
    "JOB_SLA": f"{sla_map.get('Job SLA', '')}%",
    "CYCLE_SLA_VALUE": f"{sla_map.get('Cycle SLA', '')}%"
}

# -----------------------------
# 3. LOAD PPTX TEMPLATE
# -----------------------------
prs = Presentation("template.pptx")
slide = prs.slides[0]

# -----------------------------
# 4. TEXT REPLACEMENT
# -----------------------------
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for key, value in replacements.items():
                if key in run.text:
                    run.text = run.text.replace(key, value)

# -----------------------------
# CENTER ALIGN ALL TEXT
# -----------------------------
for shape in slide.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full_text = para.text.replace("–", "-").replace("—", "-").strip()

            # ❌ MTTR description → NOT centered
            if "MTTR-Mean Time To Resolve" in full_text:
                para.alignment = PP_ALIGN.LEFT
            # ✅ Everything else → centered
            else:
                para.alignment = PP_ALIGN.CENTER

# -----------------------------
# SET FONT TO POPPINS
# -----------------------------
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue

    for para in shape.text_frame.paragraphs:
        full_text = para.text.replace("–", "-").replace("—", "-").strip()

        for run in para.runs:
            txt = run.text.strip()
            run.font.name = "Poppins"

            # Dashboard Title
            if "Operational Support" in txt:
                run.font.size = Pt(25)
                run.font.bold = True

            # ✅ ONLY this full label becomes small
            elif "MTTR-Mean Time To Resolve" in full_text:
                run.font.size = Pt(9)
                run.font.bold = True

            # KPI Values
            elif txt.replace("<", "").replace("%", "").replace(" ", "").isdigit():
                run.font.size = Pt(20)
                run.font.bold = True

            # SLA Values
            elif "%" in txt:
                run.font.size = Pt(14)
                run.font.bold = True

            # Labels
            else:
                run.font.size = Pt(13)
                run.font.bold = True

# -----------------------------
# FIX MTTR (<5) CENTER ALIGNMENT
# -----------------------------
for shape in slide.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()

        if text == "<5" or text == "< 5":
            tf = shape.text_frame
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name = "Poppins"
                    run.font.size = Pt(24)
                    run.font.bold = True

# -----------------------------
# HELPER FUNCTION: FIND CHART
# -----------------------------
def get_chart_by_title(slide, title_text):
    for shape in slide.shapes:
        if shape.has_chart:
            try:
                title = shape.chart.chart_title.text_frame.text.strip()
                if title_text.lower() in title.lower() or title.lower() in title_text.lower():
                    return shape.chart
            except AttributeError:
                continue
    return None

# -----------------------------
# 5. CHART POPULATION
# -----------------------------

# --- Chart 1: Ticket Trend ---
chart = get_chart_by_title(slide, "Ticket Count by Month")
if chart:
    fiscal_order_full = [
        "June", "July", "August", "September", "October", "November",
        "December", "January", "February", "March", "April", "May"
    ]
    month_map = {
        "Jun": "June", "Jul": "July", "Aug": "August",
        "Sep": "September", "Oct": "October", "Nov": "November",
        "Dec": "December", "Jan": "January", "Feb": "February",
        "Mar": "March", "Apr": "April", "May": "May"
    }
    reverse_map = {v: k for k, v in month_map.items()}

    trend_df['Month'] = trend_df['Month'].astype(str).str.strip()
    trend_df['MonthFull'] = trend_df['Month'].replace(month_map)
    trend_df['MonthFull'] = pd.Categorical(trend_df['MonthFull'], categories=fiscal_order_full, ordered=True)
    trend_df = trend_df.sort_values('MonthFull')
    trend_df['MonthShort'] = trend_df['MonthFull'].map(reverse_map)

    data = CategoryChartData()
    data.categories = trend_df['MonthShort']
    data.add_series('Created', trend_df['Created'])
    data.add_series('Closed', trend_df['Closed'])
    chart.replace_data(data)

# --- Chart 2: Source Availability ---
chart = get_chart_by_title(slide, "Source Data Availability")
if chart:
    data = CategoryChartData()
    data.categories = source_df['SourceSystem']
    data.add_series('On-time', source_df['On-time'])
    data.add_series('Delayed', source_df['Delayed'])
    chart.replace_data(data)

# --- Chart 3: Issue Type ---
chart = get_chart_by_title(slide, "Production Failures by Issue Type")
if chart:
    categories = issue_df['Source'].tolist()
    job_failure = [None if x == 0 else x for x in issue_df['Job Failure']]
    platform_issue = [None if x == 0 else x for x in issue_df['Platform Issue']]
    source_dq = [None if x == 0 else x for x in issue_df['Source DQ']]

    data = CategoryChartData()
    data.categories = categories
    data.add_series('Job Failure', job_failure)
    data.add_series('Platform Issue', platform_issue)
    data.add_series('Source DQ', source_dq)
    chart.replace_data(data)

# --- Chart 4: Batch Type (Pie) ---
chart = get_chart_by_title(slide, "Production Failures by Batch Type")
if chart:
    total = batch_df['Value'].sum()
    batch_df['LabelText'] = batch_df.apply(
        lambda row: f"{int(row['Value'])}, {round((row['Value']/total)*100)}%", axis=1
    )

    data = CategoryChartData()
    data.categories = batch_df['BatchType']
    data.add_series('Batches', batch_df['Value'])
    chart.replace_data(data)

    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.show_value = False
    labels.show_percentage = False
    labels.show_category_name = False

    for i, point in enumerate(plot.series[0].points):
        point.has_data_label = True
        point.data_label.text_frame.text = batch_df['LabelText'].iloc[i]
        point.data_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# --- Chart 5: DQ Issues (Sorted Descending) ---
chart = get_chart_by_title(slide, "Open Source Data Quality Issues")
if chart is None:
    chart = get_chart_by_title(slide, "Open-Source Data Quality Issue")

if chart:
    dq_df['OpenIssues'] = pd.to_numeric(dq_df['OpenIssues'])
    dq_df = dq_df.sort_values(by='OpenIssues', ascending=True)
    
    data = CategoryChartData()
    data.categories = dq_df['Source']
    data.add_series('Open Issues', dq_df['OpenIssues'])
    chart.replace_data(data)

# -----------------------------
# 6. DATA VALIDATION
# -----------------------------
if trend_df.empty or source_df.empty or issue_df.empty or batch_df.empty or dq_df.empty:
    raise ValueError("❌ Error: One or more dataframes are empty. Validation failed.")
print("✅ Validation Passed")


# -----------------------------
# AUTOMATED GAUGE GENERATION
# -----------------------------
# Extract active percentage strings directly from your data frame maps
job_sla_value = sla_map.get('Job SLA', 100)
cycle_sla_value = sla_map.get('Cycle SLA', 100)

print(f"🔄 Re-generating dynamic gauges from Excel data... Job SLA: {job_sla_value}%, Cycle SLA: {cycle_sla_value}%")

# Execute scripts seamlessly, passing the live metrics down as terminal variables
subprocess.run(["python", "job_gauge.py", str(job_sla_value)])
subprocess.run(["python", "cycle_gauge.py", str(cycle_sla_value)])


# -----------------------------
# 7. INSERT GAUGE IMAGES (TOP CENTER)
# -----------------------------
# Dimensions
gauge_width = Inches(2.65)
gauge_height = Inches(1.65)
gap = Inches(0.95)

# Calculate centered coordinates
slide_width = prs.slide_width
total_width = gauge_width * 2 + gap
start_x = Inches(4.35)
top_pos = Inches(0.85)

# ---------------- Job Gauge ----------------
job_pic = slide.shapes.add_picture(
    "job gauge.png",
    left=start_x,
    top=top_pos,
    width=gauge_width,
    height=gauge_height
)
job_pic.line.color.rgb = RGBColor(225, 225, 225)
job_pic.line.width = Pt(1)

# ---------------- Cycle Gauge ----------------
cycle_pic = slide.shapes.add_picture(
    "cycle gauge.png",
    left=start_x + gauge_width + gap,
    top=top_pos,
    width=gauge_width,
    height=gauge_height
)
cycle_pic.line.color.rgb = RGBColor(225, 225, 225)
cycle_pic.line.width = Pt(1)

# -----------------------------
# 8. SAVE FINAL DASHBOARD
# -----------------------------
output_file = "output_dashboard.pptx"
prs.save(output_file)

print(f"🎯 Complete dashboard successfully generated: {output_file}")